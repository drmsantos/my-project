import os
import csv
from pptx import Presentation
from pptx.util import Cm

def ler_configuracoes_arquivo(arquivo_config):
    configuracoes = []
    with open(arquivo_config, 'r', encoding='utf-8') as f:
        leitor_csv = csv.DictReader(f, delimiter="|")
        campos_esperados = ["diretorio", "imagem", "pag.", "largura_cm", "altura_cm", "posicao_horizontal_cm", "posicao_vertical_cm"]
        if leitor_csv.fieldnames != campos_esperados:
            raise ValueError(f"Os campos esperados são: {', '.join(campos_esperados)}. "
                             f"Campos encontrados: {', '.join(leitor_csv.fieldnames or [])}")

        for linha in leitor_csv:
            try:
                configuracoes.append({
                    "diretorio": linha["diretorio"],
                    "imagem": linha["imagem"],
                    "pagina": int(linha["pag."]),
                    "largura_cm": float(linha["largura_cm"].replace(',', '.')),
                    "altura_cm": float(linha["altura_cm"].replace(',', '.')),
                    "posicao_horizontal_cm": float(linha["posicao_horizontal_cm"].replace(',', '.')),
                    "posicao_vertical_cm": float(linha["posicao_vertical_cm"].replace(',', '.'))
                })
            except ValueError as e:
                print(f"Erro ao processar linha: {linha}. Detalhes: {e}")
    return configuracoes

def ajustar_imagens_pptx(arquivo_entrada, arquivo_saida, configuracoes, diretorio_base):
    prs = Presentation(arquivo_entrada)

    for config in configuracoes:
        slide_index = config["pagina"] - 1
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            imagem_caminho = os.path.join(diretorio_base, config["diretorio"], config["imagem"])
            if os.path.exists(imagem_caminho):
                largura = Cm(config["largura_cm"])
                altura = Cm(config["altura_cm"])
                pos_h = Cm(config["posicao_horizontal_cm"])
                pos_v = Cm(config["posicao_vertical_cm"])

                img_shape = slide.shapes.add_picture(imagem_caminho, pos_h, pos_v, largura, altura)
                slide.shapes._spTree.remove(img_shape._element)
                slide.shapes._spTree.insert(2, img_shape._element)
            else:
                print(f"Imagem não encontrada: {imagem_caminho}")
        else:
            print(f"Slide {config['pagina']} não existe.")
    prs.save(arquivo_saida)
