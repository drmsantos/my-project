import os
import csv
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Cm

# Diretório base onde as imagens podem ser encontradas
DIRETORIO_BASE = r"C:\APP\Coletas\Coleta_ElasticSearch\Imagens"

def ler_configuracoes_arquivo(arquivo_config):
    """Lê o arquivo de configuração e retorna as informações como lista de dicionários."""
    configuracoes = []
    try:
        with open(arquivo_config, 'r') as f:
            leitor_csv = csv.DictReader(f, delimiter="|")
            campos_esperados = ["diretorio", "imagem", "pag.", "largura_cm", "altura_cm", "posicao_horizontal_cm", "posicao_vertical_cm"]
            if leitor_csv.fieldnames != campos_esperados:
                raise ValueError(f"Os campos esperados são: {', '.join(campos_esperados)}. "
                                 f"Campos encontrados: {', '.join(leitor_csv.fieldnames or [])}")
            
            for linha in leitor_csv:
                try:
                    configuracoes.append({
                        "diretorio": linha["diretorio"] if linha["diretorio"] else DIRETORIO_BASE,
                        "imagem": linha["imagem"],
                        "pagina": int(linha["pag."]),
                        "largura_cm": float(linha["largura_cm"]),
                        "altura_cm": float(linha["altura_cm"]),
                        "posicao_horizontal_cm": float(linha["posicao_horizontal_cm"].replace(',', '.')),
                        "posicao_vertical_cm": float(linha["posicao_vertical_cm"].replace(',', '.'))
                    })
                except ValueError as e:
                    print(f"Erro ao processar linha: {linha}. Detalhes: {e}")
        return configuracoes
    except Exception as e:
        raise ValueError(f"Erro ao ler o arquivo de configuração: {e}")

def ajustar_imagens_pptx(arquivo_entrada, arquivo_saida, configuracoes):
    """Ajusta imagens nos slides especificados de uma apresentação PowerPoint."""
    try:
        prs = Presentation(arquivo_entrada)

        for config in configuracoes:
            slide_index = config["pagina"] - 1  # Os slides são baseados em índice 0
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]

                imagem_caminho = os.path.join(DIRETORIO_BASE, config["diretorio"], config["imagem"])

                print(f"Procurando imagem em: {imagem_caminho}")

                if os.path.exists(imagem_caminho):
                    largura = Cm(config["largura_cm"])
                    altura = Cm(config["altura_cm"])
                    posicao_horizontal = Cm(config["posicao_horizontal_cm"])
                    posicao_vertical = Cm(config["posicao_vertical_cm"])

                    imagem_shape = slide.shapes.add_picture(
                        imagem_caminho,
                        left=posicao_horizontal,
                        top=posicao_vertical,
                        width=largura,
                        height=altura
                    )

                    # Envia a imagem para trás
                    slide.shapes._spTree.remove(imagem_shape._element)
                    slide.shapes._spTree.insert(2, imagem_shape._element)

                    print(f"Imagem adicionada e enviada para trás: {imagem_caminho} no slide {config['pagina']}")
                else:
                    print(f"Imagem não encontrada: {imagem_caminho}")
            else:
                print(f"Slide {config['pagina']} não existe na apresentação.")

        prs.save(arquivo_saida)
        messagebox.showinfo("Sucesso", f"Configurações aplicadas e arquivo salvo como:\n{arquivo_saida}")
    except Exception as e:
        messagebox.showerror("Erro", f"Ocorreu um erro ao ajustar as imagens:\n{e}")

def selecionar_arquivo_entrada():
    """Abre janela para selecionar o arquivo PowerPoint de entrada."""
    return filedialog.askopenfilename(
        title="Selecione o arquivo PowerPoint",
        filetypes=[("Arquivos PowerPoint", "*.pptx")]
    )

def selecionar_arquivo_configuracoes():
    """Abre janela para selecionar o arquivo de configuração."""
    return filedialog.askopenfilename(
        title="Selecione o arquivo de configuração",
        filetypes=[("Arquivos de Texto", "*.txt")]
    )

def selecionar_diretorio_saida():
    """Abre janela para selecionar o diretório de saída."""
    return filedialog.askdirectory(title="Selecione o diretório para salvar o arquivo ajustado")

def solicitar_nome_arquivo_saida(diretorio_saida, arquivo_entrada, configuracoes):
    """Solicita o nome do arquivo de saída e executa o ajuste das imagens."""
    janela_saida = tk.Toplevel()
    janela_saida.title("Nome do Arquivo de Saída")

    tk.Label(janela_saida, text="Digite o nome do arquivo de saída:").pack(pady=10)
    nome_arquivo_entry = tk.Entry(janela_saida, width=40)
    nome_arquivo_entry.pack(pady=10)

    # Sugestão automática de nome baseado no arquivo de entrada
    nome_sugerido = os.path.splitext(os.path.basename(arquivo_entrada))[0] + "_ajustado"
    nome_arquivo_entry.insert(0, nome_sugerido)

    def salvar_nome():
        nome_arquivo = nome_arquivo_entry.get().strip()
        if nome_arquivo:
            caminho_saida = os.path.join(diretorio_saida, f"{nome_arquivo}.pptx")
            ajustar_imagens_pptx(arquivo_entrada, caminho_saida, configuracoes)
            janela_saida.destroy()
        else:
            messagebox.showwarning("Atenção", "Por favor, insira um nome válido para o arquivo.")

    tk.Button(janela_saida, text="Salvar", command=salvar_nome).pack(pady=10)

    janela_saida.protocol("WM_DELETE_WINDOW", lambda: janela_saida.destroy())

    # Evita travamento: bloqueia interação até janela ser fechada
    janela_saida.grab_set()
    janela_saida.wait_window()

if __name__ == "__main__":
    root = tk.Tk()
    root.withdraw()  # Oculta a janela principal

    arquivo_config = selecionar_arquivo_configuracoes()
    if not arquivo_config:
        messagebox.showwarning("Atenção", "Nenhum arquivo de configuração selecionado. O programa será encerrado.")
    else:
        try:
            configuracoes = ler_configuracoes_arquivo(arquivo_config)
            if not configuracoes:
                messagebox.showwarning("Atenção", "Nenhuma configuração válida encontrada. O programa será encerrado.")
            else:
                arquivo_entrada = selecionar_arquivo_entrada()
                if not arquivo_entrada:
                    messagebox.showwarning("Atenção", "Nenhum arquivo PowerPoint selecionado. O programa será encerrado.")
                else:
                    diretorio_saida = selecionar_diretorio_saida()
                    if not diretorio_saida:
                        messagebox.showwarning("Atenção", "Nenhum diretório selecionado. O programa será encerrado.")
                    else:
                        solicitar_nome_arquivo_saida(diretorio_saida, arquivo_entrada, configuracoes)
        except ValueError as e:
            messagebox.showerror("Erro", str(e))

    root.destroy()  # Finaliza corretamente a aplicação
