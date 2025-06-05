from flask import Flask, send_file, abort
import io
from pptx import Presentation

app = Flask(__name__)

@app.route('/gerar_pptx')
def gerar_pptx():
    try:
        # Cria um PPTX básico na memória
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # layout em branco
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Exemplo PPTX gerado"

        # Salva PPTX em bytes
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        # Envia arquivo para download
        return send_file(
            pptx_io,
            as_attachment=True,
            download_name='output.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        print(f"Erro ao gerar PPTX: {e}")
        return abort(500, description="Erro interno ao gerar PPTX")
