# Vizor - Aplicação Flask Corrigida
import os
import sys
import csv
import threading
import subprocess
from flask import (
    Flask, render_template, jsonify, request,
    send_file, redirect, url_for, flash
)
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Cm

app = Flask(__name__)
app.secret_key = "chave_secreta_do_vizor"

# Pasta para uploads
UPLOAD_FOLDER = os.path.join(os.getcwd(), "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Caminho fixo para imagens
CAMINHO_IMAGENS = r"/home/operador/my-project/capturas"

# Logs e controle de execução
logs = []
logs_lock = threading.Lock()
executando = threading.Event()

# Scripts de OpenSearch
caminhos_scripts = {
    'barra': '/home/operador/my-project/Api-BarOpenSearch.py',
    'funcionario': '/home/operador/my-project/Api-FunOpenSearch.py',
    'palmeiras': '/home/operador/my-project/Api-PalOpenSearch.py',
    'jaguare': '/home/operador/my-project/Api-JagOpenSearch.py'
}

# Scripts de Elastic
caminhos_elastic = {
    'barelastic': '/home/operador/my-project/Api-BarElastic.py',
    'funelastic': '/home/operador/my-project/Api-FunElastic.py',
    'palelastic': '/home/operador/my-project/Api-PalElastic.py',
    'jagelastic': '/home/operador/my-project/Api-JagElastic.py'
}

def executar_script(caminho_script):
    """Executa script e armazena logs."""
    global logs
    with logs_lock:
        logs.clear()
        logs.append("Início da execução...\n")
    try:
        process = subprocess.Popen(
            [sys.executable, caminho_script],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True
        )
        for line in iter(process.stdout.readline, ''):
            if line:
                with logs_lock:
                    logs.append(line)
        process.stdout.close()
        process.wait()
        with logs_lock:
            status = "sucesso" if process.returncode == 0 else f"erro (código {process.returncode})"
            logs.append(f"Execução finalizada com {status}.\n")
    except Exception as e:
        with logs_lock:
            logs.append(f"Erro ao executar: {e}\n")
    finally:
        executando.clear()

def ler_configuracoes_arquivo(arquivo_config):
    """Lê configurações do arquivo TXT com delimitador '|'."""
    configuracoes = []
    with open(arquivo_config, 'r', encoding='utf-8') as f:
        leitor_csv = csv.DictReader(f, delimiter='|')
        if leitor_csv.fieldnames is None:
            raise ValueError("Arquivo de configuração vazio ou mal formatado.")

        # Normaliza colunas: minúsculas e remove espaços
        colunas = [c.strip().lower() for c in leitor_csv.fieldnames]

        # Campos esperados sem underscore para comparação
        campos_esperados_sem_underscore = [
            "diretorio", "imagem", "pag", "largura_cm",
            "altura_cm", "posicaohorizontalcm", "posicaoverticalcm"
        ]

        # Mapeia colunas do arquivo para seus nomes originais
        colunas_map = {c.replace("_", ""): c for c in colunas}

        faltantes = [c for c in campos_esperados_sem_underscore if c not in colunas_map]
        if faltantes:
            raise ValueError(f"Campos ausentes: {', '.join(faltantes)}")

        for linha in leitor_csv:
            try:
                def valor(campo_sem_underscore):
                    chave = colunas_map[campo_sem_underscore]
                    return linha[chave]

                configuracoes.append({
                    "diretorio": valor("diretorio"),
                    "imagem": valor("imagem"),
                    "pagina": int(valor("pag")),
                    "largura_cm": float(valor("largura_cm").replace(',', '.')),
                    "altura_cm": float(valor("altura_cm").replace(',', '.')),
                    "posicao_horizontal_cm": float(valor("posicaohorizontalcm").replace(',', '.')),
                    "posicao_vertical_cm": float(valor("posicaoverticalcm").replace(',', '.'))
                })
            except Exception as e:
                print(f"Erro ao ler linha: {linha} -> {e}", file=sys.stderr)
    return configuracoes

def ajustar_imagens_pptx(pptx_in, pptx_out, configuracoes, base_dir):
    """Insere imagens nos slides de acordo com a configuração."""
    prs = Presentation(pptx_in)
    for config in configuracoes:
        idx = config["pagina"] - 1
        if 0 <= idx < len(prs.slides):
            slide = prs.slides[idx]
            img_path = os.path.join(base_dir, config["diretorio"], config["imagem"])
            if os.path.exists(img_path):
                shape = slide.shapes.add_picture(
                    img_path,
                    Cm(config["posicao_horizontal_cm"]),
                    Cm(config["posicao_vertical_cm"]),
                    Cm(config["largura_cm"]),
                    Cm(config["altura_cm"])
                )
                # Colocar imagem na ordem correta
                slide.shapes._spTree.remove(shape._element)
                slide.shapes._spTree.insert(2, shape._element)
            else:
                print(f"Imagem não encontrada: {img_path}", file=sys.stderr)
        else:
            print(f"Slide {config['pagina']} não existe.", file=sys.stderr)
    prs.save(pptx_out)

# ROTAS FLASK

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/elastic')
def elastic():
    return render_template('elastic.html')

@app.route('/opensearch')
def opensearch():
    return render_template('opensearch.html')

@app.route('/posdrop')
def posdrop():
    return render_template('posdrop.html')

@app.route('/sair')
def sair():
    return "Sessão finalizada. Pode fechar o navegador."

@app.route('/executar_captura', methods=['POST'])
def executar_captura():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"status": "Requisição inválida."}), 400

    nome = data.get('script')
    caminho = caminhos_scripts.get(nome)
    if not caminho:
        return jsonify({"status": f"Script '{nome}' não encontrado."}), 404

    if executando.is_set():
        return jsonify({"status": "Já existe uma execução em andamento."}), 409

    executando.set()
    threading.Thread(target=executar_script, args=(caminho,), daemon=True).start()
    return jsonify({"status": f"Execução do script '{nome}' iniciada."})

@app.route('/executar_elastic', methods=['POST'])
def executar_elastic():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"status": "Requisição inválida."}), 400

    nome = data.get('script')
    caminho = caminhos_elastic.get(nome)
    if not caminho:
        return jsonify({"status": f"Script '{nome}' não encontrado."}), 404

    if executando.is_set():
        return jsonify({"status": "Já existe uma execução em andamento."}), 409

    executando.set()
    threading.Thread(target=executar_script, args=(caminho,), daemon=True).start()
    return jsonify({"status": f"Execução do script Elastic '{nome}' iniciada."})

@app.route('/logs')
def obter_logs():
    with logs_lock:
        return jsonify({"logs": logs.copy()})

@app.route('/limpar_logs', methods=['POST'])
def limpar_logs():
    with logs_lock:
        logs.clear()
    return jsonify({"status": "Logs limpos com sucesso."})

@app.route('/pptx', methods=['GET', 'POST'])
def pptx():
    if request.method == 'POST':
        pptx_file = request.files.get('pptx_file')
        config_file = request.files.get('config_file')
        nome_saida = request.form.get('nome_saida')

        if not pptx_file or not config_file or not nome_saida:
            flash("Todos os campos são obrigatórios.")
            return redirect(url_for('pptx'))

        pptx_path = os.path.join(UPLOAD_FOLDER, secure_filename(pptx_file.filename))
        config_path = os.path.join(UPLOAD_FOLDER, secure_filename(config_file.filename))
        pptx_file.save(pptx_path)
        config_file.save(config_path)

        try:
            configuracoes = ler_configuracoes_arquivo(config_path)
            if not configuracoes:
                flash("Arquivo de configuração vazio.")
                return redirect(url_for('pptx'))

            nome_saida = secure_filename(nome_saida)
            if not nome_saida.endswith('.pptx'):
                nome_saida += '.pptx'

            saida_path = os.path.join(UPLOAD_FOLDER, nome_saida)
            ajustar_imagens_pptx(pptx_path, saida_path, configuracoes, CAMINHO_IMAGENS)
            return render_template('pptx_resultado.html', arquivo=nome_saida)
        except Exception as e:
            flash(f"Erro ao processar: {e}")
            return redirect(url_for('pptx'))

    return render_template('pptx_index.html')

@app.route('/download/<arquivo>')
def download(arquivo):
    caminho = os.path.join(UPLOAD_FOLDER, arquivo)
    if not os.path.exists(caminho):
        flash("Arquivo não encontrado.")
        return redirect(url_for('pptx'))
    return send_file(caminho, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0')
